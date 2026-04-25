"""Build a single PPT slide with all 40 WS photos in an 8-col x 5-row grid.
Slide is landscape US Letter (11" x 8.5"). Each photo has a black caption bar
below it ("WS## / location"). The user edits and exports the slide manually."""

import csv
import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).parent.parent.parent  # nus-campus-weather-paper/
PHOTO_DIR = ROOT / "ws-photos"
CSV_PATH = ROOT.parent / "BEAM_WS_Locations.csv"
TMP_DIR = ROOT.parent / "_ws_cropped_tmp"
OUT_PPTX = ROOT / "ws_stations_row.pptx"

TMP_DIR.mkdir(exist_ok=True)

# --- load locations ---
locations = {}  # id (int) -> location string
with open(CSV_PATH, "r", encoding="utf-8", errors="replace") as f:
    reader = csv.DictReader(f)
    for row in reader:
        sid = int(row["ID"])
        loc = row["Location"].strip().replace("\xa0", " ")
        locations[sid] = loc

assert len(locations) == 40, f"expected 40 stations, got {len(locations)}"

# --- center-crop each image to 3:4 portrait and save to tmp ---
TARGET_AR = 1.0  # width / height; square crop to maximize image size in the fixed slide
cropped_paths = {}
for i in range(1, 41):
    # find source file (WS18 is .jpeg, others .jpg)
    src = None
    for ext in (".jpg", ".jpeg", ".JPG", ".JPEG"):
        p = PHOTO_DIR / f"WS{i:02d}{ext}"
        if p.exists():
            src = p
            break
    assert src is not None, f"missing photo WS{i:02d}"

    im = Image.open(src)
    im = im.convert("RGB")
    # auto-rotate per EXIF
    try:
        from PIL import ImageOps
        im = ImageOps.exif_transpose(im)
    except Exception:
        pass
    w, h = im.size
    cur_ar = w / h
    if cur_ar > TARGET_AR:
        # too wide -> crop left/right
        new_w = int(h * TARGET_AR)
        left = (w - new_w) // 2
        im = im.crop((left, 0, left + new_w, h))
    else:
        # too tall -> crop top/bottom
        new_h = int(w / TARGET_AR)
        top = (h - new_h) // 2
        im = im.crop((0, top, w, top + new_h))
    # resize to a sensible size to keep PPT filesize reasonable
    im = im.resize((800, 800), Image.LANCZOS)
    out = TMP_DIR / f"WS{i:02d}.jpg"
    im.save(out, "JPEG", quality=85)
    cropped_paths[i] = out

# --- build PPT ---
SLIDE_W_IN = 11.0   # landscape US Letter
SLIDE_H_IN = 8.5
prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

N_COLS = 8
N_ROWS = 5
IMG_W_IN = 1.25
IMG_H_IN = 1.25         # 1:1 square; fills both slide axes
BAR_H_IN = 0.32
COL_GAP_IN = 0.08
ROW_GAP_IN = 0.06

IMG_W = Inches(IMG_W_IN)
IMG_H = Inches(IMG_H_IN)
BAR_H = Inches(BAR_H_IN)
COL_GAP = Inches(COL_GAP_IN)
ROW_GAP = Inches(ROW_GAP_IN)

# center grid on the slide
grid_w_in = N_COLS * IMG_W_IN + (N_COLS - 1) * COL_GAP_IN
grid_h_in = N_ROWS * (IMG_H_IN + BAR_H_IN) + (N_ROWS - 1) * ROW_GAP_IN
LEFT0 = Inches((SLIDE_W_IN - grid_w_in) / 2)
TOP0 = Inches((SLIDE_H_IN - grid_h_in) / 2)
CELL_H = IMG_H + BAR_H

for i in range(1, 41):
    idx = i - 1
    row, col = idx // N_COLS, idx % N_COLS
    left = LEFT0 + (IMG_W + COL_GAP) * col
    top = TOP0 + (CELL_H + ROW_GAP) * row

    # image
    slide.shapes.add_picture(str(cropped_paths[i]), left, top, width=IMG_W, height=IMG_H)

    # black caption bar below
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + IMG_H, IMG_W, BAR_H)
    bar.line.fill.background()  # no border
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 0, 0)

    tf = bar.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # line 1: WS##
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = f"WS{i:02d}"
    r1.font.bold = True
    r1.font.size = Pt(8)
    r1.font.color.rgb = RGBColor(255, 255, 255)

    # line 2: location
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = locations[i]
    r2.font.size = Pt(6)
    r2.font.color.rgb = RGBColor(255, 255, 255)

# ============================================================
# Slide 2: 4-variable annual summary table
# Day = 07:00-18:59 SGT; Night = 19:00-06:59 SGT
# ============================================================
import pandas as pd

DATA_DIR = ROOT / "data" / "imputed"

VAR_COL = {
    "AirTemp":   "AirTemp Ave (C)",
    "RelHum":    "RelHum Ave (%)",
    "GlobalRad": "GlobalRad Ave (W/m2)",
    "WindSpeed": "WindSpeed Ave (m/s)",
}

# (row label, source var, stat, time-of-day filter or None, round decimals)
# None entries are blank separator rows inserted between variable sections.
TABLE_ROWS = [
    ("AirTemp Max (°C)",        "AirTemp",    "max",  None,    2),
    ("AirTemp AvgDay (°C)",     "AirTemp",    "mean", "day",   2),
    ("AirTemp Min (°C)",        "AirTemp",    "min",  None,    2),
    ("AirTemp AvgNight (°C)",   "AirTemp",    "mean", "night", 2),
    ("__BLANK__",                    None,         None,   None,    None),
    ("RH Max (%)",                   "RelHum",     "max",  None,    1),
    ("RH Avg (%)",                   "RelHum",     "mean", None,    1),
    ("RH Min (%)",                   "RelHum",     "min",  None,    1),
    ("__BLANK__",                    None,         None,   None,    None),
    ("GlobalRad Max (W/m²)",    "GlobalRad",  "max",  None,    1),
    ("GlobalRad Avg (W/m²)",    "GlobalRad",  "mean", None,    1),
    ("__BLANK__",                    None,         None,   None,    None),
    ("WindSpeed AvgDay (m/s)",       "WindSpeed",  "mean", "day",   2),
    ("WindSpeed AvgNight (m/s)",     "WindSpeed",  "mean", "night", 2),
]

# compute values per station
station_values = {}  # i -> list aligned with TABLE_ROWS
for i in range(1, 41):
    fp = DATA_DIR / f"NUS_CAMPUS_WS{i:02d}_2025_Hourly_imputed.csv"
    df = pd.read_csv(fp, parse_dates=["Datetime"])
    hrs = df["Datetime"].dt.hour
    day_mask = hrs.between(7, 18)
    night_mask = (hrs < 7) | (hrs > 18)
    vals = []
    for label, var, stat, tod, decimals in TABLE_ROWS:
        if var is None:
            vals.append("")
            continue
        s = df[VAR_COL[var]]
        if tod == "day":
            s = s[day_mask]
        elif tod == "night":
            s = s[night_mask]
        v = getattr(s, stat)()
        vals.append(round(float(v), decimals))
    station_values[i] = vals

# add a new blank slide
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

N_TROWS = 1 + len(TABLE_ROWS)   # 1 header + 14 body = 15 rows
N_TCOLS = 1 + 40                # 1 label col + 40 WS cols

TBL_LEFT = Inches(0.25)
TBL_TOP = Inches(0.60)
TBL_W = Inches(SLIDE_W_IN - 0.50)
TBL_H = Inches(SLIDE_H_IN - 0.90)

tbl_shape = slide2.shapes.add_table(N_TROWS, N_TCOLS, TBL_LEFT, TBL_TOP, TBL_W, TBL_H)
tbl = tbl_shape.table

# column widths: first column wider for labels
label_w = Inches(1.55)
data_w = Inches((SLIDE_W_IN - 0.50 - 1.55) / 40)
tbl.columns[0].width = label_w
for c in range(1, N_TCOLS):
    tbl.columns[c].width = data_w

# header row
hdr = tbl.rows[0]
hdr.cells[0].text = "Metric"
for c in range(1, N_TCOLS):
    hdr.cells[c].text = f"WS{c:02d}"

def _fmt_cell(cell, text, *, bold=False, size=6, fill=None, align=PP_ALIGN.CENTER):
    cell.text = ""  # reset to strip default paragraph
    tf = cell.text_frame
    tf.margin_left = Emu(18000)
    tf.margin_right = Emu(18000)
    tf.margin_top = Emu(9000)
    tf.margin_bottom = Emu(9000)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill

# style header
for c in range(N_TCOLS):
    cell = tbl.rows[0].cells[c]
    txt = cell.text
    _fmt_cell(cell, txt, bold=True, size=7, fill=RGBColor(0x1F, 0x3A, 0x5F))
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)

# body rows
for r_idx, (label, var, stat, tod, decimals) in enumerate(TABLE_ROWS, start=1):
    row = tbl.rows[r_idx]
    is_blank = (var is None)
    row.height = Inches(0.18 if is_blank else 0.38)

    if is_blank:
        # leave all cells empty but strip default bullet paragraphs
        for c in range(N_TCOLS):
            _fmt_cell(row.cells[c], "", size=6)
        continue

    # label cell
    _fmt_cell(row.cells[0], label, bold=True, size=7,
              fill=RGBColor(0xEA, 0xEA, 0xEA), align=PP_ALIGN.LEFT)
    # data cells
    for i in range(1, 41):
        v = station_values[i][r_idx - 1]
        text = f"{v}" if v != "" else ""
        _fmt_cell(row.cells[i], text, size=6)

prs.save(OUT_PPTX)
print(f"saved: {OUT_PPTX}")
print(f"slide 1: {SLIDE_W_IN} x {SLIDE_H_IN} in photo grid ({N_COLS}x{N_ROWS})")
print(f"slide 2: {N_TROWS}x{N_TCOLS} table, {len(TABLE_ROWS)} body rows (3 blank separators)")
